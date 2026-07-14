from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation object
prs = Presentation()

# Size of slide 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define theme color (dark mode premium look)
BG_COLOR = RGBColor(18, 30, 38)
TEXT_CYAN = RGBColor(0, 212, 255)
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_MUTED = RGBColor(93, 173, 226)

# Slide database (with new slides and expanded information)
slides_data = [
    {
        "type": "title",
        "title": "Advanced Digital Clock &\nDesktop Utility Suite",
        "subtitle": "A Professional, Multi-Featured Python Application with Cross-Platform Utilities\n\nPresented by: MD. Atiqur Rahman, Tajim & Tarek\nCourse/Project Code: Project Works-1(28556)"
    },
    {
        "type": "content",
        "title": "Project Overview & Objectives",
        "points": [
            "Transforming a simple digital clock into a powerful desktop utility tool.",
            "Integrating productivity modules like Pomodoro, Timer, and Stopwatch into a single UI.",
            "Creating a rich World Clock database featuring 100+ global and commercial cities.",
            "Monitoring real-time system resources (Battery & Charging Status).",
            "Enhancing User Experience (UX) with dynamic themes, opacity slider, and custom fonts."
        ]
    },
    {
        "type": "content",
        "title": "Problem Statement & Motivation",
        "points": [
            "Traditional OS clocks only show time and require external, heavy apps for timers or alarms.",
            "Lack of integrated focus tools (like Pomodoro) forces users to switch tabs, breaking deep work flow.",
            "Most desktop widgets consume high RAM and lack cross-platform transparency controls.",
            "Motivation: Build a resource-friendly (<50MB RAM), centralized desktop assistant using pure Python."
        ]
    },
    {
        "type": "content",
        "title": "Core Architecture & Tech Stack",
        "points": [
            "Programming Language: Python 3.x - Lightweight and highly efficient.",
            "GUI Framework: Tkinter & Ttk - Used for professional and structured widget styling.",
            "Hardware Sync: psutil - Fetches real-time laptop battery percentage and status.",
            "Desktop Notifications: plyer - Handles cross-platform system push notifications.",
            "Timezone Engine: zoneinfo - Implements the IANA timezone database for precise global time."
        ]
    },
    {
        "type": "content",
        "title": "Key Features: Time & Productivity",
        "points": [
            "Smart Alarm Manager: Allows setting, managing, and deleting multiple alarms smoothly.",
            "Stopwatch Mode: High-precision millisecond tracking with dedicated start, pause, and reset options.",
            "Countdown Timer: User-friendly combo-box inputs for precise time tracking with alerts.",
            "Pomodoro Tracker: Automated 25-minute work and 5-minute break sessions to boost efficiency."
        ]
    },
    {
        "type": "content",
        "title": "Advanced Features: Global Sync & Reminders",
        "points": [
            "Massive World Clock Database: 100 major cities categorized by continent to avoid UI clutter.",
            "Event & Birthday Reminder: In-memory database that alerts users on specific dates upon app startup.",
            "Hourly Chime System: Triggers an OS-level push notification and audio alert (bell) at every hour completion."
        ]
    },
    {
        "type": "content",
        "title": "UI/UX & Customization Features",
        "points": [
            "Dynamic Theme Customizer: 4 built-in professional themes (Standard Dark, Light, Neon Green, Neon Pink).",
            "Transparency Controls: Adjusts window opacity (0.3 to 1.0) for better screen visibility during multitask.",
            "Font Style Selector: One-click font swapping (Arial, Times New Roman, Courier New, Impact) across the app."
        ]
    },
    {
        "type": "content",
        "title": "Special Utilities for Power Users",
        "points": [
            "Compact / Mini Mode: Resizes the massive UI into a sleek 300x80 widget to save desktop space.",
            "Battery Status Monitor: Embedded directly in the footer, showing charging state dynamically.",
            "Floating Mode (Always on Top): Keeps the clock pinned above all other running applications."
        ]
    },
    {
        "type": "content",
        "title": "Application Workflow (How it Works)",
        "points": [
            "Initialization: App loads the main Tkinter thread, fetches local time, and links the IANA timezone database.",
            "Asynchronous Looping: The update_time() function runs non-blocking tasks every 1000ms using .after().",
            "Background Checking: Simultaneously scans active alarms, checks hourly chime logic, and requests battery telemetry.",
            "UI Interactivity: Menu items trigger independent Toplevel windows for features like Stopwatch or Reminders."
        ]
    },
    {
        "type": "content",
        "title": "Comparison: Built-in OS Clock vs Our Solution",
        "points": [
            "Feature Depth: OS Clock shows basic time | Our App provides Stopwatch, Timer, Pomodoro, and Reminders.",
            "Screen Real Estate: OS Clock is fixed on the taskbar | Our App features Floating Mode and a scalable Mini View.",
            "Personalization: Minimal OS theme choices | Our App supports 4 neon/light themes and runtime opacity scaling.",
            "Resource Overhead: Heavy modern widgets use high memory | Our Python core is ultra-lightweight and standalone."
        ]
    },
    {
        "type": "content",
        "title": "Technical Challenges & Solutions",
        "points": [
            "UI Freezing Issue: Fixed by handling asynchronous time loops via Tkinter's .after() method instead of standard infinite loops.",
            "Timezone & DST Management: Solved by utilizing 'zoneinfo' instead of manual hardcoded offsets.",
            "Dynamic Window Resizing: Managed smoothly using precise widget layout manipulation (pack_forget and place)."
        ]
    },
    {
        "type": "content",
        "title": "Future Enhancements",
        "points": [
            "Database Integration: Implementing SQLite for data persistence (saving alarms and reminders permanently).",
            "Cloud Synchronization: Syncing data with Google Calendar API and external task managers.",
            "Custom Audio Uploads: Allowing users to set custom .mp3/.wav files as alarm tones.",
            "System Tray App: Minimizing the software directly into the OS system tray tree."
        ]
    },
    {
        "type": "title",
        "title": "Thank You",
        "subtitle": "Open for Questions & Feedback\n\nEmail: contact.atiqur.tajim.tarek@example.com\nGitHub: github.com/yourusername/digital-clock"
    }
]

# Slide creation loop
for data in slides_data:
    # Taking a blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Setting the background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    if data["type"] == "title":
        # Title slide formatting
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.33), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = data["title"]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TEXT_CYAN
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = data["subtitle"]
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_WHITE
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(24)
        
    elif data["type"] == "content":
        # Content slide header
        titleBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1))
        tf_title = titleBox.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = data["title"]
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_CYAN
        
        # Content body (Bulit point)
        contentBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5))
        tf_content = contentBox.text_frame
        tf_content.word_wrap = True
        
        for i, point in enumerate(data["points"]):
            if i == 0:
                p_point = tf_content.paragraphs[0]
            else:
                p_point = tf_content.add_paragraph()
            
            p_point.text = "• " + point
            p_point.font.size = Pt(18)
            p_point.font.color.rgb = TEXT_WHITE
            p_point.space_after = Pt(14)
            p_point.level = 0

# save the file
prs.save("Digital_Clock_Advanced_Presentation1.pptx")
print("Presentation file 'Digital_Clock_Advanced_Presentation.pptx' has been created successfully!")
